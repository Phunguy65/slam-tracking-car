#!/usr/bin/env python3
"""Export presentation.html as editable PPTX using python-pptx."""

import os
import sys
from pathlib import Path
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

CANVAS_W = Inches(13.333)
CANVAS_H = Inches(7.5)
MARGIN_X = Inches(0.6)
MARGIN_TOP = Inches(0.5)
CONTENT_MAX_Y = Inches(6.70)
FOOTER_TOP = Inches(6.85)
FOOTER_H = Inches(0.22)

COLOR_INK_BG = RGBColor(0x0E, 0x11, 0x17)
COLOR_LIGHT_BG = RGBColor(0xFA, 0xF9, 0xF7)
COLOR_INK_FG = RGBColor(0xF0, 0xED, 0xE8)
COLOR_LIGHT_FG = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_MUTED_LIGHT = RGBColor(0x5C, 0x5C, 0x5C)
COLOR_MUTED_INK = RGBColor(0x88, 0x88, 0x88)
COLOR_ACCENT_LIGHT = RGBColor(0x2C, 0x4A, 0x6E)
COLOR_ACCENT_INK = RGBColor(0x7E, 0xB8, 0xDA)

FONT_DISPLAY = "Georgia"
FONT_BODY = "Arial"

CONTENT_W = CANVAS_W - 2 * MARGIN_X


class Cursor:
    def __init__(self, y_start=MARGIN_TOP, cap=CONTENT_MAX_Y):
        self.y = y_start
        self.cap = cap

    def take(self, h, gap=Inches(0.12)):
        top = self.y
        self.y = top + h + gap
        return top

    def remaining(self):
        return self.cap - self.y


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=Pt(14), bold=False, italic=False, color=None,
                alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_footer(slide, left_text, right_text, is_ink=False):
    fg = COLOR_MUTED_INK if is_ink else COLOR_MUTED_LIGHT
    add_textbox(slide, MARGIN_X, FOOTER_TOP, Inches(6), FOOTER_H,
                left_text, font_size=Pt(11), color=fg)
    add_textbox(slide, CANVAS_W - MARGIN_X - Inches(4), FOOTER_TOP,
                Inches(4), FOOTER_H, right_text, font_size=Pt(11),
                color=fg, alignment=PP_ALIGN.RIGHT)


def add_header(slide, left, center, right, is_ink=False):
    fg = COLOR_MUTED_INK if is_ink else COLOR_MUTED_LIGHT
    top = Inches(0.25)
    h = Inches(0.25)
    add_textbox(slide, MARGIN_X, top, Inches(3), h, left,
                font_size=Pt(10), color=fg)
    add_textbox(slide, Inches(5.5), top, Inches(3), h, center,
                font_size=Pt(10), color=fg, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, CANVAS_W - MARGIN_X - Inches(2), top, Inches(2), h,
                right, font_size=Pt(10), color=fg, alignment=PP_ALIGN.RIGHT)


def add_kicker(slide, top, text, is_ink=False):
    color = COLOR_ACCENT_INK if is_ink else COLOR_ACCENT_LIGHT
    add_textbox(slide, MARGIN_X, top, CONTENT_W, Inches(0.25), text,
                font_size=Pt(11), bold=True, color=color)


def add_title(slide, top, text, is_ink=False, font_size=Pt(36)):
    color = COLOR_INK_FG if is_ink else COLOR_LIGHT_FG
    add_textbox(slide, MARGIN_X, top, CONTENT_W, Inches(0.8), text,
                font_name=FONT_DISPLAY, font_size=font_size, color=color)


def build_slide_01(prs):
    """Trang bìa"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_INK_BG)
    c = Cursor(Inches(0.8))

    add_textbox(slide, MARGIN_X, c.take(Inches(0.2)),
                Inches(8), Inches(0.25),
                "Học viện Công nghệ Bưu chính Viễn thông",
                font_size=Pt(11), color=COLOR_MUTED_INK)

    add_kicker(slide, c.take(Inches(0.25)), "INT1461 — Xây dựng các hệ thống nhúng", is_ink=True)

    add_textbox(slide, MARGIN_X, c.take(Inches(1.2)),
                Inches(9), Inches(1.2),
                "Hệ thống Xe Robot\nTự hành SLAM",
                font_name=FONT_DISPLAY, font_size=Pt(54), color=COLOR_INK_FG)

    add_textbox(slide, MARGIN_X, c.take(Inches(0.4)),
                Inches(8), Inches(0.4),
                "Khám phá bản đồ • Điều hướng tự động • Theo dõi người",
                font_size=Pt(18), color=COLOR_MUTED_INK)

    add_textbox(slide, MARGIN_X, c.take(Inches(0.3)),
                Inches(6), Inches(0.3),
                "Giảng viên hướng dẫn: Nguyễn Trọng Huân",
                font_size=Pt(13), color=COLOR_MUTED_INK)

    members_top = Inches(1.2)
    members_left = Inches(9.0)
    add_textbox(slide, members_left, members_top, Inches(3.8), Inches(0.25),
                "Thành viên nhóm 2", font_size=Pt(10), color=COLOR_MUTED_INK)
    members = [
        ("Nguyễn Ngọc Phú", "N22DCCN159"),
        ("Ngô Tấn Sang", "N22DCCN167"),
        ("Văn Minh Tấn", "N22DCCN175"),
        ("Vũ Tiến Đạt", "N22DCCN120"),
        ("Huỳnh Phát Tài", "N22DCCN171"),
    ]
    for i, (name, code) in enumerate(members):
        y = members_top + Inches(0.4 + i * 0.4)
        add_textbox(slide, members_left, y, Inches(2.2), Inches(0.35),
                    name, font_size=Pt(16), color=COLOR_INK_FG)
        add_textbox(slide, members_left + Inches(2.2), y, Inches(1.2), Inches(0.35),
                    code, font_size=Pt(11), color=COLOR_MUTED_INK, alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, members_left, members_top + Inches(2.5), Inches(3.8), Inches(0.25),
                "TP.HCM, tháng 06 / 2025", font_size=Pt(11), color=COLOR_MUTED_INK)

    add_footer(slide, "SLAM Tracking Car", "Báo cáo đồ án", is_ink=True)


def build_slide_02(prs):
    """Tổng quan hệ thống — Ba chế độ hoạt động"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "SLAM Tracking Car", "Tổng quan", "02 / 16")
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Kiến trúc hệ thống")
    add_title(slide, c.take(Inches(0.7)), "Ba chế độ hoạt động")

    modes = [
        ("01", "Khám phá bản đồ", "SLAM Toolbox xây dựng bản đồ từ LiDAR + explore_lite tự động di chuyển đến vùng chưa khám phá"),
        ("02", "Điều hướng tự động", "Nav2 Stack: AMCL định vị + NavFn lập đường đi + DWB điều khiển bám đường"),
        ("03", "Theo dõi người", "YOLOv8n + InsightFace nhận dạng, Camera-LiDAR fusion định vị, PID 3 tầng điều khiển"),
    ]
    col_w = Inches(3.8)
    top = c.take(Inches(2.5))
    for i, (num, title, desc) in enumerate(modes):
        left = MARGIN_X + i * (col_w + Inches(0.3))
        add_textbox(slide, left, top, col_w, Inches(0.5), num,
                    font_name=FONT_DISPLAY, font_size=Pt(30), color=COLOR_LIGHT_FG)
        add_textbox(slide, left, top + Inches(0.6), col_w, Inches(0.4), title,
                    font_name=FONT_DISPLAY, font_size=Pt(18), color=COLOR_LIGHT_FG)
        add_textbox(slide, left, top + Inches(1.1), col_w, Inches(1.2), desc,
                    font_size=Pt(14), color=COLOR_MUTED_LIGHT)

    flow_top = c.take(Inches(1.0))
    flow_text = "LiDAR LDS02RR → SLAM Toolbox → Nav2 Stack → Motor Driver    |    ESP32-CAM"
    add_textbox(slide, MARGIN_X, flow_top, CONTENT_W, Inches(0.4),
                flow_text, font_size=Pt(14), color=COLOR_MUTED_LIGHT,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide, "Nguồn: Tài liệu thiết kế hệ thống", "Nhóm 2 — 2026")


def build_slide_03(prs):
    """Kiến trúc phần cứng — with robot image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_INK_BG)
    add_header(slide, "SLAM Tracking Car", "Phần cứng", "03 / 16", is_ink=True)
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Thiết kế phần cứng", is_ink=True)
    add_title(slide, c.take(Inches(0.7)), "Kiến trúc khối hệ thống", is_ink=True)

    body_top = c.take(Inches(4.5))
    img_path = Path(__file__).parent / "assets" / "robot.jpg"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), MARGIN_X, body_top,
                                 width=Inches(4.5), height=Inches(4.0))
    add_textbox(slide, MARGIN_X, body_top + Inches(4.1), Inches(4.5), Inches(0.25),
                "Mô hình xe robot thực tế", font_size=Pt(11),
                color=COLOR_MUTED_INK, alignment=PP_ALIGN.CENTER)

    blocks_left = MARGIN_X + Inches(5.2)
    blocks_w = Inches(3.5)
    blocks = [
        ("Khối Cảm biến", ["LiDAR LDS02RR — UART 360°", "MPU6050 IMU — I2C 6-DoF",
                           "Encoder × 2 — Interrupt", "ESP32-CAM — HTTP MJPEG"]),
        ("Khối Xử lý", ["ESP32 Main — micro-ROS", "ESP32-CAM — Streamer", "WiFi UDP → Agent"]),
        ("Khối Chấp hành", ["TB6612 — PWM 2 kênh", "2× Servo SG90 — Pan/Tilt"]),
        ("Khối Nguồn", ["3 cặp Li-ion 18650", "2× Buck LM2596", "Pin trực tiếp → Motor"]),
    ]
    for idx, (title, items) in enumerate(blocks):
        col = idx % 2
        row = idx // 2
        left = blocks_left + col * (blocks_w + Inches(0.3))
        top = body_top + row * Inches(2.0)
        add_textbox(slide, left, top, blocks_w, Inches(0.3), title,
                    font_size=Pt(14), bold=True, color=COLOR_ACCENT_INK)
        item_text = "\n".join(f"• {it}" for it in items)
        add_textbox(slide, left, top + Inches(0.35), blocks_w, Inches(1.5),
                    item_text, font_size=Pt(13), color=COLOR_INK_FG)

    cost_top = body_top + Inches(4.1)
    add_textbox(slide, blocks_left, cost_top, Inches(3.5), Inches(0.3),
                "Tổng chi phí phần cứng", font_size=Pt(12), color=COLOR_MUTED_INK)
    add_textbox(slide, blocks_left + Inches(3.8), cost_top, Inches(3), Inches(0.3),
                "~1.200.000 VNĐ", font_name=FONT_DISPLAY, font_size=Pt(20),
                color=COLOR_INK_FG, alignment=PP_ALIGN.RIGHT)

    add_footer(slide, "Thiết kế module hóa — dễ khắc phục sự cố và nâng cấp",
               "Nhóm 2 — 2026", is_ink=True)


def build_slide_04(prs):
    """Firmware module hóa trên ESP32"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "SLAM Tracking Car", "Firmware", "04 / 16")
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "PlatformIO + micro-ROS")
    add_title(slide, c.take(Inches(0.7)), "Firmware module hóa trên ESP32")

    body_top = c.take(Inches(4.5))
    left_col = MARGIN_X
    col_w = Inches(5.8)

    add_textbox(slide, left_col, body_top, col_w, Inches(0.35),
                "ESP32 Main — micro-ROS Node",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)

    table_data = [
        ["Module", "Publish", "Subscribe"],
        ["motors", "—", "/cmd_vel"],
        ["encoders", "/odom", "—"],
        ["imu", "/imu/data_raw", "—"],
        ["lidar", "/scan", "—"],
        ["servos", "/joint_states", "/servo_cmd"],
        ["safety", "Watchdog: dừng motor khi mất agent", ""],
        ["ros_bridge", "Interface chung các module", ""],
    ]
    tbl_top = body_top + Inches(0.5)
    rows, cols = len(table_data), 3
    table_shape = slide.shapes.add_table(rows, cols, left_col, tbl_top,
                                         col_w, Inches(2.8))
    table = table_shape.table
    for r, row_data in enumerate(table_data):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(r, ci)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.name = FONT_BODY

    right_col = MARGIN_X + Inches(6.5)
    right_w = Inches(5.5)
    add_textbox(slide, right_col, body_top, right_w, Inches(0.35),
                "ESP32-CAM — Video Streamer",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)
    cam_text = "Phát luồng MJPEG qua HTTP :80\nKhông phải micro-ROS node\ncam_bridge_node kéo dữ liệu về ROS"
    add_textbox(slide, right_col, body_top + Inches(0.5), right_w, Inches(1.0),
                cam_text, font_size=Pt(14), color=COLOR_MUTED_LIGHT)

    add_textbox(slide, right_col, body_top + Inches(1.8), right_w, Inches(0.35),
                "Cơ chế an toàn",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)
    safety_items = "• Watchdog timer — dừng motor khi mất kết nối\n• PID anti-windup trên 2 kênh motor\n• WiFi reconnect tự động\n• Cấu hình .env: WiFi, IP agent, GPIO pins"
    add_textbox(slide, right_col, body_top + Inches(2.2), right_w, Inches(1.5),
                safety_items, font_size=Pt(14), color=COLOR_LIGHT_FG)

    add_footer(slide, "Mỗi ngoại vi có file source riêng — giao tiếp qua ros_bridge interface",
               "Nhóm 2 — 2026")


def build_slide_05(prs):
    """SLAM — Scan Matching"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "SLAM Tracking Car", "Khám phá bản đồ", "05 / 16")
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "SLAM Toolbox")
    add_title(slide, c.take(Inches(0.7)), "So khớp quét tương quan")

    body_top = c.take(Inches(4.5))
    left_col = MARGIN_X
    col_w = Inches(5.8)

    desc = "Ước lượng phép biến đổi giữa 2 bản quét LiDAR liên tiếp bằng phương pháp correlative scan matching kết hợp tối ưu phi tuyến Ceres Solver."
    add_textbox(slide, left_col, body_top, col_w, Inches(1.0),
                desc, font_size=Pt(15), color=COLOR_MUTED_LIGHT)

    add_textbox(slide, left_col, body_top + Inches(1.2), col_w, Inches(0.4),
                "E(T) = Σᵢ [1 − M(T · pᵢ)]²",
                font_name="Times New Roman", font_size=Pt(20),
                italic=True, color=COLOR_LIGHT_FG)

    params = "• T — phép biến đổi cần tìm (dx, dy, dθ)\n• pᵢ — điểm thứ i trong bản quét hiện tại\n• M(x) — xác suất tại vị trí x trên bản đồ"
    add_textbox(slide, left_col, body_top + Inches(1.8), col_w, Inches(1.2),
                params, font_size=Pt(15), color=COLOR_LIGHT_FG)

    right_col = MARGIN_X + Inches(6.5)
    right_w = Inches(5.5)
    steps = "Bước 1: Xây dựng bản đồ xác suất\n     ↓\nBước 2: Quét thô (bước 5cm, 2°)\n  Tìm T cho tổng xác suất cao nhất\n     ↓\nBước 3: Tối ưu tinh — Levenberg-Marquardt\n  Hội tụ đến vị trí chính xác sub-pixel\n     ↓\n  → Phép biến đổi T tối ưu"
    add_textbox(slide, right_col, body_top, right_w, Inches(4.0),
                steps, font_size=Pt(14), color=COLOR_LIGHT_FG)

    add_footer(slide, "Thuật toán: SPARSE_NORMAL_CHOLESKY + SCHUR_JACOBI",
               "Nhóm 2 — 2026")


def build_slide_06(prs):
    """Loop Closure + Occupancy Grid"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "SLAM Tracking Car", "Khám phá bản đồ", "06 / 16")
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Đóng vòng lặp & Bản đồ lưới")
    add_title(slide, c.take(Inches(0.7)), "Loop Closure + Occupancy Grid")

    body_top = c.take(Inches(4.5))
    left_col = MARGIN_X
    col_w = Inches(5.8)

    add_textbox(slide, left_col, body_top, col_w, Inches(0.35),
                "Phát hiện vòng lặp",
                font_name=FONT_DISPLAY, font_size=Pt(18), color=COLOR_LIGHT_FG)
    desc = "Khi robot quay lại vùng đã khám phá, thuật toán phát hiện trùng lặp và hiệu chỉnh sai số tích lũy qua tối ưu toàn cục đồ thị tư thế."
    add_textbox(slide, left_col, body_top + Inches(0.5), col_w, Inches(0.8),
                desc, font_size=Pt(14), color=COLOR_MUTED_LIGHT)

    table_data = [
        ["Tham số", "Giá trị"],
        ["Chuỗi scan tối thiểu", "10"],
        ["Ngưỡng matching thô", "≥ 0.35"],
        ["Ngưỡng matching tinh", "≥ 0.45"],
        ["Variance tối đa (thô)", "≤ 3.0"],
    ]
    tbl_top = body_top + Inches(1.5)
    tbl = slide.shapes.add_table(len(table_data), 2, left_col, tbl_top,
                                 col_w, Inches(2.0))
    for r, row_data in enumerate(table_data):
        for ci, txt in enumerate(row_data):
            tbl.table.cell(r, ci).text = txt

    right_col = MARGIN_X + Inches(6.5)
    right_w = Inches(5.5)
    add_textbox(slide, right_col, body_top, right_w, Inches(0.35),
                "Bản đồ Occupancy Grid",
                font_name=FONT_DISPLAY, font_size=Pt(18), color=COLOR_LIGHT_FG)
    grid_desc = "Trắng = Trống (0)\nĐen = Vật cản (100)\nXám = Chưa biết (−1)\nXanh = Robot\nĐộ phân giải: 0.05m/pixel\n\nCập nhật mỗi 5 giây. Mỗi ô chứa xác suất chiếm đóng từ 0 (trống) đến 100 (vật cản)."
    add_textbox(slide, right_col, body_top + Inches(0.5), right_w, Inches(3.0),
                grid_desc, font_size=Pt(14), color=COLOR_MUTED_LIGHT)

    add_footer(slide, "Cấu hình: map_update_interval = 5.0s", "Nhóm 2 — 2026")


def build_slide_07(prs):
    """Frontier Exploration"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "SLAM Tracking Car", "Khám phá bản đồ", "07 / 16")
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "explore_lite")
    add_title(slide, c.take(Inches(0.7)), "Khám phá dựa trên biên giới")

    body_top = c.take(Inches(4.5))
    left_col = MARGIN_X
    col_w = Inches(5.8)

    desc = "Robot tự động tìm và di chuyển đến các biên giới — ranh giới giữa vùng đã biết và chưa biết — để khám phá toàn bộ môi trường."
    add_textbox(slide, left_col, body_top, col_w, Inches(0.8),
                desc, font_size=Pt(15), color=COLOR_MUTED_LIGHT)

    add_textbox(slide, left_col, body_top + Inches(1.0), col_w, Inches(0.4),
                "score(f) = α × size(f) − β × cost(f) − γ × turn(f)",
                font_name="Times New Roman", font_size=Pt(18),
                italic=True, color=COLOR_LIGHT_FG)

    table_data = [
        ["Trọng số", "Giá trị", "Ý nghĩa"],
        ["α (gain_scale)", "0.8", "Ưu tiên biên giới lớn"],
        ["β (potential_scale)", "5.0", "Ưu tiên biên giới gần"],
        ["γ (orientation_scale)", "0.0", "Bỏ qua góc quay"],
    ]
    tbl_top = body_top + Inches(1.6)
    tbl = slide.shapes.add_table(len(table_data), 3, left_col, tbl_top,
                                 col_w, Inches(1.6))
    for r, row_data in enumerate(table_data):
        for ci, txt in enumerate(row_data):
            tbl.table.cell(r, ci).text = txt

    right_col = MARGIN_X + Inches(6.5)
    right_w = Inches(5.5)
    steps = "Quét bản đồ hiện tại\n     ↓\nTìm tất cả biên giới (BFS)\n     ↓\nChấm điểm và chọn biên giới tốt nhất\n     ↓\nGửi mục tiêu đến Nav2\n\n↻ Lặp mỗi 3s"
    add_textbox(slide, right_col, body_top, right_w, Inches(3.5),
                steps, font_size=Pt(15), color=COLOR_LIGHT_FG)

    add_footer(slide, "Điều kiện dừng: không còn biên giới ≥ 0.35m", "Nhóm 2 — 2026")


def build_slide_08(prs):
    """Phối hợp SLAM + Exploration"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_INK_BG)
    add_header(slide, "SLAM Tracking Car", "Khám phá bản đồ", "08 / 16", is_ink=True)
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Luồng dữ liệu", is_ink=True)
    add_title(slide, c.take(Inches(0.7)), "Phối hợp SLAM Toolbox & explore_lite", is_ink=True)

    body_top = c.take(Inches(4.0))
    flow = "LiDAR → /scan → SLAM Toolbox → /map → explore_lite → goal → Nav2 Stack → /cmd_vel → Motor"
    add_textbox(slide, MARGIN_X, body_top + Inches(1.5), CONTENT_W, Inches(0.5),
                flow, font_size=Pt(16), color=COLOR_INK_FG, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN_X, body_top + Inches(2.5), CONTENT_W, Inches(0.4),
                "Robot di chuyển → LiDAR quét vùng mới → vòng lặp tiếp tục",
                font_size=Pt(14), italic=True, color=COLOR_MUTED_INK,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide, "Vòng lặp khám phá tự động hoàn toàn", "Nhóm 2 — 2026", is_ink=True)


def build_slide_09(prs):
    """Nav2 — Kiến trúc phân tầng"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "SLAM Tracking Car", "Điều hướng", "09 / 16")
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Nav2 Stack")
    add_title(slide, c.take(Inches(0.7)), "Kiến trúc phân tầng điều hướng")

    body_top = c.take(Inches(4.5))

    add_textbox(slide, Inches(4.5), body_top, Inches(4.5), Inches(0.5),
                "BT Navigator (Điều phối tổng thể)",
                font_size=Pt(14), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.CENTER)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(4.5), body_top, Inches(4.5), Inches(0.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_ACCENT_LIGHT

    servers = [
        ("Planner Server", "Lập đường đi toàn cục (A*)"),
        ("Controller Server", "Bám đường cục bộ (DWB)"),
        ("Behavior Server", "Phục hồi khi kẹt"),
    ]
    for i, (name, desc) in enumerate(servers):
        left = MARGIN_X + Inches(0.3) + i * Inches(4.1)
        top = body_top + Inches(1.0)
        add_textbox(slide, left, top, Inches(3.6), Inches(0.3), name,
                    font_size=Pt(13), bold=True, color=COLOR_LIGHT_FG)
        add_textbox(slide, left, top + Inches(0.35), Inches(3.6), Inches(0.3),
                    desc, font_size=Pt(12), color=COLOR_MUTED_LIGHT)

    costmaps = [
        "Global Costmap (Static + Obstacle + Inflation)",
        "Local Costmap (Rolling window)",
        "Spin / Backup / Wait",
    ]
    for i, cm in enumerate(costmaps):
        left = MARGIN_X + Inches(0.3) + i * Inches(4.1)
        top = body_top + Inches(2.0)
        add_textbox(slide, left, top, Inches(3.6), Inches(0.3), cm,
                    font_size=Pt(12), color=COLOR_MUTED_LIGHT)

    lifecycle = "Thứ tự kích hoạt Lifecycle: map_server → amcl → controller → planner → bt_navigator"
    add_textbox(slide, MARGIN_X, body_top + Inches(3.2), CONTENT_W, Inches(0.3),
                lifecycle, font_size=Pt(12), color=COLOR_MUTED_LIGHT,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide, "Quản lý bởi MapManagerNode (ROS2 Lifecycle)", "Nhóm 2 — 2026")


def build_slide_10(prs):
    """AMCL — Bộ lọc hạt thích nghi"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "SLAM Tracking Car", "Điều hướng", "10 / 16")
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Định vị trên bản đồ")
    add_title(slide, c.take(Inches(0.7)), "AMCL — Bộ lọc hạt thích nghi")

    body_top = c.take(Inches(4.5))
    left_col = MARGIN_X
    col_w = Inches(5.8)

    desc = "Sử dụng 500–2000 hạt (particles) để ước lượng vị trí robot. Số lượng hạt tự điều chỉnh dựa trên KL-divergence."
    add_textbox(slide, left_col, body_top, col_w, Inches(0.7),
                desc, font_size=Pt(14), color=COLOR_MUTED_LIGHT)

    steps = [
        ("1. Khởi tạo", "N hạt phân bố Gaussian quanh vị trí ban đầu"),
        ("2. Dự đoán", "Di chuyển hạt theo odometry + nhiễu (α₁–α₅ = 0.2)"),
        ("3. Cập nhật", "Tính trọng số bằng likelihood field model"),
        ("4. Lấy mẫu lại", "Loại hạt trọng số thấp, nhân bản hạt trọng số cao"),
    ]
    for i, (title, desc_text) in enumerate(steps):
        top = body_top + Inches(0.9 + i * 0.7)
        add_textbox(slide, left_col, top, col_w, Inches(0.25), title,
                    font_size=Pt(13), bold=True, color=COLOR_ACCENT_LIGHT)
        add_textbox(slide, left_col, top + Inches(0.28), col_w, Inches(0.3),
                    desc_text, font_size=Pt(12), color=COLOR_MUTED_LIGHT)

    right_col = MARGIN_X + Inches(6.5)
    right_w = Inches(5.5)
    table_data = [
        ["Tham số", "Giá trị"],
        ["Số hạt tối thiểu", "500"],
        ["Số hạt tối đa", "2000"],
        ["Mô hình cảm biến", "Likelihood field"],
        ["Số tia laser sử dụng", "180 / 360"],
    ]
    tbl = slide.shapes.add_table(len(table_data), 2, right_col,
                                 body_top + Inches(2.0), right_w, Inches(2.0))
    for r, row_data in enumerate(table_data):
        for ci, txt in enumerate(row_data):
            tbl.table.cell(r, ci).text = txt

    add_footer(slide, "Mô hình: nav2_amcl::DifferentialMotionModel", "Nhóm 2 — 2026")


def build_slide_11(prs):
    """NavFn (A*) + DWB"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "SLAM Tracking Car", "Điều hướng", "11 / 16")
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Lập đường đi & Bám đường")
    add_title(slide, c.take(Inches(0.7)), "NavFn (A*) + Dynamic Window (DWB)")

    body_top = c.take(Inches(4.5))
    left_col = MARGIN_X
    col_w = Inches(5.8)

    add_textbox(slide, left_col, body_top, col_w, Inches(0.3),
                "Thuật toán A*",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)
    add_textbox(slide, left_col, body_top + Inches(0.4), col_w, Inches(0.35),
                "f(n) = g(n) + h(n)",
                font_name="Times New Roman", font_size=Pt(18), italic=True,
                color=COLOR_LIGHT_FG)
    astar_params = "• g(n) — chi phí thực từ start đến n\n• h(n) — heuristic Euclidean đến goal\n• Dung sai goal: 0.5m"
    add_textbox(slide, left_col, body_top + Inches(0.9), col_w, Inches(0.9),
                astar_params, font_size=Pt(14), color=COLOR_LIGHT_FG)

    add_textbox(slide, left_col, body_top + Inches(2.0), col_w, Inches(0.3),
                "Costmap Inflation",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)
    add_textbox(slide, left_col, body_top + Inches(2.4), col_w, Inches(0.35),
                "cost(d) = 253 × e^(−10 × (d − 0.12))",
                font_name="Times New Roman", font_size=Pt(16), italic=True,
                color=COLOR_LIGHT_FG)
    add_textbox(slide, left_col, body_top + Inches(2.9), col_w, Inches(0.3),
                "Vùng đệm 45cm quanh vật cản, bán kính robot 12cm",
                font_size=Pt(13), color=COLOR_MUTED_LIGHT)

    right_col = MARGIN_X + Inches(6.5)
    right_w = Inches(5.5)
    add_textbox(slide, right_col, body_top, right_w, Inches(0.3),
                "Dynamic Window Approach",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)
    add_textbox(slide, right_col, body_top + Inches(0.4), right_w, Inches(0.6),
                "Mô phỏng 1600 quỹ đạo (40×40 mẫu) trong 1.2 giây, chấm điểm bằng 7 critics.",
                font_size=Pt(14), color=COLOR_MUTED_LIGHT)

    table_data = [
        ["Critic", "Scale"],
        ["RotateToGoal", "32.0"],
        ["PathAlign", "20.0"],
        ["PathDist", "20.0"],
        ["GoalAlign", "16.0"],
        ["GoalDist", "16.0"],
        ["BaseObstacle", "0.05"],
    ]
    tbl = slide.shapes.add_table(len(table_data), 2, right_col,
                                 body_top + Inches(1.2), right_w, Inches(2.5))
    for r, row_data in enumerate(table_data):
        for ci, txt in enumerate(row_data):
            tbl.table.cell(r, ci).text = txt

    add_textbox(slide, right_col, body_top + Inches(3.9), right_w, Inches(0.3),
                "Tốc độ tối đa: 0.3 m/s tịnh tiến, 1.0 rad/s quay",
                font_size=Pt(13), color=COLOR_LIGHT_FG)

    add_footer(slide, "Goal tolerance: xy=15cm, yaw=0.25rad", "Nhóm 2 — 2026")


def build_slide_12(prs):
    """Tracking — Kiến trúc 3 tầng"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_INK_BG)
    add_header(slide, "SLAM Tracking Car", "Theo dõi người", "12 / 16", is_ink=True)
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Person Tracking", is_ink=True)
    add_title(slide, c.take(Inches(0.7)), "Kiến trúc 3 tầng theo dõi người", is_ink=True)

    body_top = c.take(Inches(4.0))

    tiers = [
        ("TẦNG 1: NHẬN DẠNG", ["ESP32-CAM (MJPEG stream)", "LiDAR LDS02RR (360° scan)",
                                 "→ Person Tracker", "  YOLOv8n (body)", "  InsightFace (face)",
                                 "  Leg Clustering"]),
        ("TẦNG 2: ĐỊNH VỊ", ["Sensor Fusion", "  Camera bearing", "  + LiDAR range",
                              "  → (range, bearing)"]),
        ("TẦNG 3: ĐIỀU KHIỂN", ["Tracking Controller", "  Servo PID (50 Hz)",
                                  "  Wheel Yaw PID (10 Hz)", "  Linear PID (10 Hz)",
                                  "  Safety Arc Check", "→ /servo_cmd, /cmd_vel"]),
    ]
    col_w = Inches(3.8)
    for i, (tier_name, items) in enumerate(tiers):
        left = MARGIN_X + i * (col_w + Inches(0.3))
        add_textbox(slide, left, body_top, col_w, Inches(0.25),
                    tier_name, font_size=Pt(11), bold=True, color=COLOR_ACCENT_INK)
        item_text = "\n".join(items)
        add_textbox(slide, left, body_top + Inches(0.4), col_w, Inches(3.0),
                    item_text, font_size=Pt(13), color=COLOR_INK_FG)

    add_textbox(slide, MARGIN_X, body_top + Inches(3.5), Inches(4), Inches(0.3),
                "Enrollment Node (SQLite database)",
                font_size=Pt(12), color=COLOR_MUTED_INK)

    add_footer(slide, "3 tầng: Nhận dạng → Định vị → Điều khiển",
               "Nhóm 2 — 2026", is_ink=True)


def build_slide_13(prs):
    """YOLO + InsightFace"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "SLAM Tracking Car", "Theo dõi người", "13 / 16")
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Phát hiện & Nhận dạng")
    add_title(slide, c.take(Inches(0.7)), "YOLOv8n + InsightFace")

    body_top = c.take(Inches(4.5))
    left_col = MARGIN_X
    col_w = Inches(5.8)

    add_textbox(slide, left_col, body_top, col_w, Inches(0.3),
                "Phát hiện cơ thể — YOLOv8n",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)
    yolo_items = "• Mô hình one-stage, chỉ detect class \"person\"\n• Ngưỡng confidence ≥ 0.5\n• Bounding box chuẩn hóa [0, 1]\n• Theo dõi liên tục qua IoU ≥ 0.3"
    add_textbox(slide, left_col, body_top + Inches(0.4), col_w, Inches(1.2),
                yolo_items, font_size=Pt(14), color=COLOR_LIGHT_FG)

    add_textbox(slide, left_col, body_top + Inches(1.8), col_w, Inches(0.3),
                "Nhận dạng khuôn mặt — InsightFace",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)
    face_items = "• Mô hình buffalo_l, vector 512 chiều\n• Chuẩn hóa L2 + cosine similarity\n• Ngưỡng nhận dạng: ≥ 0.6\n• Confidence decay: 0.1/giây khi mất mặt"
    add_textbox(slide, left_col, body_top + Inches(2.2), col_w, Inches(1.2),
                face_items, font_size=Pt(14), color=COLOR_LIGHT_FG)

    right_col = MARGIN_X + Inches(6.5)
    right_w = Inches(5.5)
    pipeline = "Ảnh từ ESP32-CAM\n     ↓\nYOLOv8n → Bounding box người\n     ↓\nCắt vùng body → InsightFace detect\n     ↓\nTrích xuất embedding 512-d\n     ↓\nSo khớp cosine ≥ 0.6 → Xác nhận"
    add_textbox(slide, right_col, body_top, right_w, Inches(4.0),
                pipeline, font_size=Pt(14), color=COLOR_LIGHT_FG)

    add_footer(slide, "Track bị xóa sau 1 giây không phát hiện", "Nhóm 2 — 2026")


def build_slide_14(prs):
    """Sensor Fusion — Camera-LiDAR"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "SLAM Tracking Car", "Theo dõi người", "14 / 16")
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Kết hợp cảm biến")
    add_title(slide, c.take(Inches(0.7)), "Camera-LiDAR Bearing Fusion")

    body_top = c.take(Inches(4.5))
    left_col = MARGIN_X
    col_w = Inches(5.8)

    add_textbox(slide, left_col, body_top, col_w, Inches(0.3),
                "Vấn đề",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)
    problem = "Camera cho biết hướng nhưng không biết khoảng cách. LiDAR biết khoảng cách nhưng không phân biệt người. Cần kết hợp 2 nguồn."
    add_textbox(slide, left_col, body_top + Inches(0.4), col_w, Inches(0.8),
                problem, font_size=Pt(14), color=COLOR_MUTED_LIGHT)

    add_textbox(slide, left_col, body_top + Inches(1.4), col_w, Inches(0.3),
                "Tính bearing từ camera",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)
    add_textbox(slide, left_col, body_top + Inches(1.8), col_w, Inches(0.35),
                "θ_cam = arctan((u − cₓ) / fₓ)",
                font_name="Times New Roman", font_size=Pt(18), italic=True,
                color=COLOR_LIGHT_FG)
    add_textbox(slide, left_col, body_top + Inches(2.3), col_w, Inches(0.3),
                "Chuyển đổi frame camera → laser qua TF2",
                font_size=Pt(13), color=COLOR_MUTED_LIGHT)

    add_textbox(slide, left_col, body_top + Inches(2.8), col_w, Inches(0.3),
                "Ghép cặp chân từ LiDAR",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)
    table_data = [
        ["Tham số", "Giá trị"],
        ["Chiều rộng chân", "5–30 cm"],
        ["Khoảng cách 2 chân", "15–35 cm"],
        ["Bearing tolerance", "0.10 rad (~5.7°)"],
    ]
    tbl = slide.shapes.add_table(len(table_data), 2, left_col,
                                 body_top + Inches(3.0), col_w, Inches(1.3))
    for r, row_data in enumerate(table_data):
        for ci, txt in enumerate(row_data):
            tbl.table.cell(r, ci).text = txt

    right_col = MARGIN_X + Inches(6.5)
    right_w = Inches(5.5)
    fusion_desc = "Kết quả fusion:\n• bearing: từ camera\n• range: từ leg pair\n\nNếu Δθ > 0.10 rad → không ghép,\ntrả range = NaN"
    add_textbox(slide, right_col, body_top + Inches(1.5), right_w, Inches(2.5),
                fusion_desc, font_size=Pt(15), color=COLOR_LIGHT_FG)

    add_footer(slide, "Nếu Δθ > 0.10 rad → không ghép, trả range = NaN",
               "Nhóm 2 — 2026")


def build_slide_15(prs):
    """PID 3 tầng + FSM"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "SLAM Tracking Car", "Theo dõi người", "15 / 16")
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Điều khiển")
    add_title(slide, c.take(Inches(0.7)), "PID 3 tầng + Máy trạng thái")

    body_top = c.take(Inches(4.5))
    left_col = MARGIN_X
    col_w = Inches(5.8)

    add_textbox(slide, left_col, body_top, col_w, Inches(0.35),
                "u(t) = Kp·e(t) + Ki·∫e(τ)dτ + Kd·de/dt",
                font_name="Times New Roman", font_size=Pt(18), italic=True,
                color=COLOR_LIGHT_FG)

    table_data = [
        ["Tầng PID", "Tần số", "Kp", "Mục đích"],
        ["Servo", "50 Hz", "2.0", "Giữ target giữa khung"],
        ["Wheel Yaw", "10 Hz", "0.5", "Quay body khi servo > 30°"],
        ["Linear", "10 Hz", "0.3", "Giữ khoảng cách 1.0–1.5m"],
    ]
    tbl = slide.shapes.add_table(len(table_data), 4, left_col,
                                 body_top + Inches(0.6), col_w, Inches(1.5))
    for r, row_data in enumerate(table_data):
        for ci, txt in enumerate(row_data):
            tbl.table.cell(r, ci).text = txt

    add_textbox(slide, left_col, body_top + Inches(2.4), col_w, Inches(0.6),
                "An toàn:\nCung quét ±20° phía trước, dừng nếu vật cản < 30cm",
                font_size=Pt(14), color=COLOR_LIGHT_FG)

    right_col = MARGIN_X + Inches(6.5)
    right_w = Inches(5.5)
    add_textbox(slide, right_col, body_top, right_w, Inches(0.3),
                "Máy trạng thái tìm kiếm",
                font_name=FONT_DISPLAY, font_size=Pt(16), color=COLOR_LIGHT_FG)
    fsm = "TRACKING\n  ↓ (mất > 0.5s)\nSEARCH_CONTINUE (0.5s)\n  ↓\nSEARCH_SCAN (2.0s)\n  ↓\nSEARCH_ROTATE (5.0s)\n  ↓\nIDLE\n\n← Target lại → quay về TRACKING"
    add_textbox(slide, right_col, body_top + Inches(0.5), right_w, Inches(3.5),
                fsm, font_size=Pt(14), color=COLOR_LIGHT_FG)

    add_footer(slide, "Target xuất hiện lại → quay về TRACKING ngay lập tức",
               "Nhóm 2 — 2026")


def build_slide_16(prs):
    """Tổng kết — Bảng tổng hợp giải thuật"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_INK_BG)
    add_header(slide, "SLAM Tracking Car", "Tổng kết", "16 / 16", is_ink=True)
    c = Cursor(Inches(1.0))
    add_kicker(slide, c.take(Inches(0.25)), "Tổng kết", is_ink=True)
    add_title(slide, c.take(Inches(0.7)), "Bảng tổng hợp giải thuật", is_ink=True)

    body_top = c.take(Inches(4.5))
    table_data = [
        ["Chức năng", "Thuật toán chính", "Thư viện"],
        ["Xây dựng bản đồ", "Graph-based SLAM, Ceres Solver", "SLAM Toolbox"],
        ["Khám phá tự động", "Frontier-based Exploration", "explore_lite"],
        ["Định vị", "Adaptive Monte Carlo Localization", "Nav2 AMCL"],
        ["Lập đường đi", "A* trên costmap", "Nav2 NavFn"],
        ["Bám đường", "Dynamic Window Approach", "Nav2 DWB"],
        ["Phát hiện người", "CNN one-stage detection", "YOLOv8n"],
        ["Nhận dạng mặt", "ArcFace + Cosine similarity", "InsightFace"],
        ["Định vị người", "Bearing fusion + Leg clustering", "Custom"],
        ["Điều khiển", "PID 3 tầng + FSM", "Custom"],
    ]
    rows, cols = len(table_data), 3
    tbl_shape = slide.shapes.add_table(rows, cols, MARGIN_X, body_top,
                                       CONTENT_W, Inches(4.0))
    table = tbl_shape.table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(5.5)
    table.columns[2].width = Inches(3.133)
    for r, row_data in enumerate(table_data):
        for ci, txt in enumerate(row_data):
            cell = table.cell(r, ci)
            cell.text = txt
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    run.font.name = FONT_BODY
                    run.font.color.rgb = COLOR_INK_FG

    add_footer(slide, "Nhóm 2 — 2026", "Cảm ơn đã lắng nghe", is_ink=True)


def main():
    prs = Presentation()
    prs.slide_width = CANVAS_W
    prs.slide_height = CANVAS_H

    build_slide_01(prs)
    build_slide_02(prs)
    build_slide_03(prs)
    build_slide_04(prs)
    build_slide_05(prs)
    build_slide_06(prs)
    build_slide_07(prs)
    build_slide_08(prs)
    build_slide_09(prs)
    build_slide_10(prs)
    build_slide_11(prs)
    build_slide_12(prs)
    build_slide_13(prs)
    build_slide_14(prs)
    build_slide_15(prs)
    build_slide_16(prs)

    output_path = Path(__file__).parent / "presentation.pptx"
    prs.save(str(output_path))
    print(f"Exported {len(prs.slides)} slides to: {output_path}")
    print(f"File size: {output_path.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
